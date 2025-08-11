import os
from pptx import Presentation
from google import genai
from google.genai import types

def setup_genai_client():
    """Setup Google GenAI client with your API key"""
    
    api_key = os.getenv("GENAI_API_KEY") 
    client = genai.Client(api_key=api_key)
    return client

def extract_slide_content(filename):
    """Extract content from PowerPoint slides"""
    print(f"üìÅ Extracting content from {filename}...")
    
    presentation = Presentation(filename)
    slides_content = []
    
    for slide_num, slide in enumerate(presentation.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        
        slide_content = '\n'.join(slide_text) if slide_text else "[No text content]"
        slides_content.append({
            'slide_number': slide_num,
            'content': slide_content
        })
        
        print(f"üìÑ Slide {slide_num}: {len(slide_text)} text elements")
    
    return slides_content

def analyze_with_genai(client, slides_content):
    """Use Google GenAI to detect inconsistencies"""
    print("\nü§ñ Analyzing with Google Gemini...")
    
    slides_text = ""
    for slide in slides_content:
        slides_text += f"\n--- SLIDE {slide['slide_number']} ---\n{slide['content']}\n"
    
    prompt = f"""
You are an expert presentation analyzer. Analyze this PowerPoint presentation for inconsistencies and conflicts.

Look for:
1. **NUMERICAL CONFLICTS**: Same metrics with different values across slides
2. **CONTRADICTORY STATEMENTS**: Claims that contradict each other  
3. **TIMELINE MISMATCHES**: Dates or sequences that don't align
4. **LOGICAL INCONSISTENCIES**: Information that doesn't make sense together

Here's the presentation content:
{slides_text}

Provide your analysis in this exact format:

ANALYSIS RESULTS:

Issues Found: [number]

Issue 1: [Type of Conflict]
- Slides: [slide numbers involved]
- Description: [clear explanation of the inconsistency]
- Details: [specific conflicting information]
- Severity: [High/Medium/Low]

Issue 2: [Type of Conflict]
[continue for each issue found]

If no conflicts are found, respond with:
ANALYSIS RESULTS:
Issues Found: 0
Status: No significant inconsistencies detected in this presentation.
"""

    try:
        response = client.models.generate_content(
            model='gemini-2.5-flash', 
            contents=prompt,
            config=types.GenerateContentConfig(
                temperature=0.1,
                max_output_tokens=2000,
            )
        )
        
        return response.text
        
    except Exception as e:
        return f"‚ùå AI Analysis failed: {str(e)}"

def main():
    print("üöÄ AI-Powered PowerPoint Inconsistency Detector")
    print("ü§ñ Using Google Gemini Flash")
    print("="*60)
    
    try:
        client = setup_genai_client()
        print("‚úÖ Google GenAI client initialized")
    except Exception as e:
        print(f"‚ùå Failed to initialize GenAI client: {e}")
        print("üí° Make sure to set your API key!")
        return
    
    filename = "test.pptx"
    
    try:
        slides_content = extract_slide_content(filename)
        print(f"‚úÖ Extracted content from {len(slides_content)} slides")
        
        ai_results = analyze_with_genai(client, slides_content)
        
        print("\n" + "="*60)
        print("üéØ INCONSISTENCY ANALYSIS RESULTS")
        print("="*60)
        print(ai_results)
        
    except FileNotFoundError:
        print(f"‚ùå File '{filename}' not found!")
    except Exception as e:
        print(f"‚ùå Error: {e}")

if __name__ == "__main__":
    main()