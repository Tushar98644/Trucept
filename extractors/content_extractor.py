from pptx import Presentation
from typing import List, Dict
from .content_handler import CONTENT_HANDLERS
from .utils import validate_file, combine_content, identify_shape_type

class ContentExtractor:    
    def __init__(self):
        self.supported_formats = ['.pptx']
    
    def extract_content(self, filename: str) -> List[Dict]:
        validate_file(filename, self.supported_formats)
        
        presentation = Presentation(filename)
        slides_content = []
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_content = self.extract_slide_content(slide, slide_num)
            slides_content.append(slide_content)
        
        return slides_content
    
    def extract_slide_content(self, slide, slide_num: int) -> Dict:
        extracted_content = []
        content_summary = {}
        
        for shape in slide.shapes:
            content_type, content_data = self.extract_from_shape(shape)
            
            if content_data:
                extracted_content.append({
                    'type': content_type,
                    'content': content_data
                })
                
                content_summary[content_type] = content_summary.get(content_type, 0) + 1
        
        combined_text = combine_content(extracted_content)
        
        return {
            'slide_number': slide_num,
            'content': combined_text,
            'content_types': content_summary,
            'total_elements': len(extracted_content),
            'raw_elements': extracted_content
        }
    
    def extract_from_shape(self, shape) -> tuple:
        content_type = identify_shape_type(shape)
        
        handler = CONTENT_HANDLERS.get(content_type)
        
        if handler:
            try:
                content_data = handler(shape)
                return content_type, content_data
            except Exception:
                return content_type, f"[Error extracting {content_type}]"
        else:
            return 'unknown', f"[{str(shape.shape_type)} shape]"

def extract_presentation(filename: str) -> List[Dict]:
    extractor = ContentExtractor()
    print(extractor.extract_content(filename))
    return extractor.extract_content(filename)